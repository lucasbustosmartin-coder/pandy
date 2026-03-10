#!/usr/bin/env python3
"""
Genera la presentación PowerPoint de la propuesta comercial (Propuesta_Pandi.pptx).
Usa el logo en ../assets/SP_logo.svg (convertido a PNG para insertar) o SP_logo.png.
La barra de título usa el mismo azul que el logo (#0d2137); el logo se genera/convierte
con fondo transparente para que luzca bien sobre la barra.
Ejecutar desde la raíz: python presentacion/crear_presentacion_pptx.py
"""
import os
import sys
from pathlib import Path

# Asegurar que podemos importar pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Instalá las dependencias: pip install -r presentacion/requirements.txt")
    sys.exit(1)

# Rutas: script puede estar en presentacion/ o invocarse desde raíz
SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
ASSETS_DIR = PROJECT_ROOT / "assets"
SVG_LOGO = ASSETS_DIR / "SP_logo.svg"
LOGO_PNG_CANDIDATES = [
    ASSETS_DIR / "SP_logo.png",   # PNG con fondo transparente recomendado
    SCRIPT_DIR / "_logo.png",     # Generado desde SVG o Pillow (con transparencia)
]
OUT_PPTX = SCRIPT_DIR / "Propuesta_Pandi.pptx"
BITACORA_XLSX = PROJECT_ROOT / "Bitacora_tareas.xlsx"

# Estilo comercial: azul del logo (L&P) y gris — fondo de barra acorde al logo, logo con transparencia
AZUL_LOGO = RGBColor(13, 33, 55)   # #0d2137 — mismo que el círculo del logo, para barra y acentos
GRIS_OSCURO = RGBColor(26, 26, 26)
GRIS_TEXTO = RGBColor(85, 85, 85)
BLANCO = RGBColor(255, 255, 255)

# Interlineado y espaciado (legibilidad)
LINE_SPACING_BODY = 1.2        # múltiplo de línea en texto cuerpo
SPACE_BEFORE_ITEM = Pt(8)      # espacio antes de cada ítem en listas
SPACE_AFTER_ITEM = Pt(2)       # espacio después de ítem
ALCANCE_ITEM_SPACE = Pt(5)     # espacio entre ítems en alcance


def leer_presupuesto_desde_excel():
    """Lee la solapa Presupuesto de Bitacora_tareas.xlsx. Devuelve lista de (funcionalidad, horas, importe_str)."""
    if not BITACORA_XLSX.exists():
        return None
    try:
        import openpyxl
        wb = openpyxl.load_workbook(BITACORA_XLSX, read_only=True, data_only=True)
        if "Presupuesto" not in wb.sheetnames:
            wb.close()
            return None
        ws = wb["Presupuesto"]
        filas = list(ws.iter_rows(min_row=1, max_col=4, values_only=True))
        wb.close()
        if len(filas) < 2:
            return None
        # filas[0] = encabezado; resto = data. Columnas 0=Funcionalidad, 1=Horas, 2=Importe, 3=Nuevo
        resultado = []
        for row in filas:
            if row[0] is None:
                continue
            func = str(row[0]).strip()
            horas = row[1]
            imp = row[2]
            if isinstance(horas, (int, float)) and not isinstance(horas, bool):
                h_str = str(int(horas))
            else:
                h_str = str(horas) if horas is not None else ""
            if isinstance(imp, (int, float)) and not isinstance(imp, bool):
                i_str = str(int(imp))
            else:
                i_str = str(imp) if imp is not None else ""
            resultado.append((func, h_str, i_str))
        return resultado if resultado else None
    except Exception:
        return None


def _generar_logo_pillow():
    """Genera un PNG del logo tipo SP (círculo oscuro + L&P) con fondo transparente. Pillow si no hay PNG/SVG usable."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None
    out = SCRIPT_DIR / "_logo.png"
    w = h = 200
    # RGBA para fondo transparente (0,0,0,0) — luce bien sobre la barra azul
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Círculo #0d2137 con contorno blanco para que resalte sobre la barra
    margin = 8
    draw.ellipse([margin, margin, w - margin, h - margin], fill=(13, 33, 55, 255), outline=(255, 255, 255, 255), width=3)
    # Texto L&P centrado en blanco
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 72)
    except Exception:
        font = ImageFont.load_default()
    text = "L&P"
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.text(((w - tw) // 2, (h - th) // 2 - 4), text, fill=(255, 255, 255, 255), font=font)
    img.save(out)
    return out


def obtener_logo_png():
    """Devuelve la ruta a un PNG del logo: primero busca PNG en assets, luego SVG→PNG, luego genera con Pillow."""
    for p in LOGO_PNG_CANDIDATES:
        if p.exists():
            return p
    if SVG_LOGO.exists():
        try:
            import cairosvg
            out = SCRIPT_DIR / "_logo.png"
            cairosvg.svg2png(url=str(SVG_LOGO), write_to=str(out), output_width=200, output_height=200)
            return out
        except Exception:
            pass
    # Fallback: generar logo tipo SP con Pillow (círculo + L&P)
    return _generar_logo_pillow()


def add_slide_number(slide, current, total):
    """Añade en la esquina inferior derecha el texto 'actual / total' (ej. 1 / 8)."""
    box = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = f"{current} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRIS_TEXTO
    p.alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, prs, titulo, logo_path=None):
    """Añade barra de título con fondo azul (logo), logo a la izquierda y título en blanco."""
    bar_h = Inches(0.75)
    # 1) Fondo verde
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, bar_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = AZUL_LOGO
    rect.line.fill.background()
    # 2) Logo a la izquierda
    if logo_path and logo_path.exists():
        logo_size = Inches(0.55)
        slide.shapes.add_picture(str(logo_path), Inches(0.2), Inches(0.1), logo_size, logo_size)
    # 3) Título en blanco a la derecha del logo
    box = slide.shapes.add_textbox(Inches(1.0), Inches(0.18), Inches(11.3), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLANCO
    return bar_h


def add_title_slide(prs, titulo, subtitulo=None, pie=None, logo_path=None):
    """Portada: barra azul con logo + título, subtítulo centrado y pie de confidencialidad."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, titulo, logo_path=logo_path)
    top = Inches(1.0)
    box2 = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(1.0))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = subtitulo or ""
    p2.font.size = Pt(18)
    p2.font.color.rgb = GRIS_TEXTO
    p2.alignment = PP_ALIGN.CENTER
    p2.line_spacing = LINE_SPACING_BODY
    if pie:
        box3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45))
        p3 = box3.text_frame.paragraphs[0]
        p3.text = pie
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRIS_TEXTO
        p3.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, titulo, lineas, es_lista=True, logo_path=None):
    """Añade slide con barra de título y contenido en viñetas o párrafos (interlineado mejorado)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, titulo, logo_path=logo_path)
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0))
    tf = content.text_frame
    tf.word_wrap = True
    for i, linea in enumerate(lineas):
        if i == 0:
            par = tf.paragraphs[0]
        else:
            par = tf.add_paragraph()
        t = linea.strip() if es_lista else linea
        if es_lista and t and not (t.startswith("•") or t.startswith("-")):
            t = "• " + t
        par.text = t
        par.font.size = Pt(14)
        par.font.color.rgb = GRIS_OSCURO
        par.line_spacing = LINE_SPACING_BODY
        if es_lista and t:
            par.space_before = SPACE_BEFORE_ITEM
            par.space_after = SPACE_AFTER_ITEM
    return slide


def add_beneficios_slide(prs, logo_path=None):
    """Slide de beneficios: barra de título + tabla 2 columnas con buen interlineado."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Qué gana tu operación", logo_path=logo_path)
    filas = [
        ("Trazabilidad", "Cada orden y transacción registrada y vinculada a caja y cuenta corriente."),
        ("Control por roles", "Admin, encargado y visor ven solo lo que necesitan; permisos por vista y por acción."),
        ("Menos errores", "Concertación única, validaciones y mensajería clara."),
        ("Múltiples monedas", "USD, EUR, ARS en cajas y cuenta corriente; tipos de operación configurables."),
        ("Uso en celular", "Interfaz adaptable para consultas y tareas básicas en movimiento."),
    ]
    rows, cols = len(filas) + 1, 2
    left, top, width = Inches(0.5), Inches(1.0), Inches(12.3)
    row_h = Inches(0.52)
    height = row_h * rows
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    table.cell(0, 0).text = "Beneficio"
    table.cell(0, 1).text = "Impacto"
    for j in range(2):
        c = table.cell(0, j)
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.color.rgb = BLANCO
        c.text_frame.paragraphs[0].font.size = Pt(12)
        c.fill.solid()
        c.fill.fore_color.rgb = AZUL_LOGO
    for i, (ben, imp) in enumerate(filas):
        table.cell(i + 1, 0).text = ben
        table.cell(i + 1, 1).text = imp
        for j in range(2):
            p = table.cell(i + 1, j).text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.line_spacing = LINE_SPACING_BODY
            p.space_before = Pt(4)
            p.space_after = Pt(4)


def add_alcance_slide(prs, logo_path=None):
    """Slide con barra de título + lista numerada de alcance (interlineado mejorado)."""
    alcance = [
        "Setup y arquitectura – Repo, Supabase, Vercel, configuración y despliegue.",
        "Autenticación y seguridad – Login, roles (Admin, Encargado, Visor), permisos y RLS.",
        "UI base – Sidebar, vistas, navegación, login y registro.",
        "Clientes – Alta, edición y listado con permisos.",
        "Cajas – Saldos por moneda, movimientos manuales y por orden, Efectivo/Banco.",
        "Órdenes – Alta, edición, estados, concertación, monedas e instrumentación.",
        "Transacciones – Ingreso/egreso, intermediarios, comisiones, tipos de operación.",
        "Cuenta corriente – Cliente e intermediario, conciliación y movimientos.",
        "Panel de Control – Saldos, pendientes, accesos (parametrizable por rol).",
        "Permisos granulares – Control por vistas y por acciones.",
        "Experiencia de usuario – Toasts, confirmaciones, validaciones.",
        "Responsive – Adaptación móvil, touch y buenas prácticas.",
        "Documentación y despliegue continuo.",
    ]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Alcance del desarrollo", logo_path=logo_path)
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.85))
    tf = content.text_frame
    tf.word_wrap = True
    for i, linea in enumerate(alcance):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.text = f"{i+1}. {linea}"
        par.font.size = Pt(12)
        par.font.color.rgb = GRIS_OSCURO
        par.line_spacing = LINE_SPACING_BODY
        par.space_before = ALCANCE_ITEM_SPACE
        par.space_after = SPACE_AFTER_ITEM


# Fallback si no hay Excel
PRESUPUESTO_FALLBACK = [
    ("Funcionalidad", "Horas", "Importe (USD)"),
    ("Setup y arquitectura", "12", "456"),
    ("Autenticación y seguridad", "24", "912"),
    ("UI base y navegación", "20", "760"),
    ("ABM Clientes", "8", "304"),
    ("Módulo Cajas", "24", "912"),
    ("Módulo Órdenes", "32", "1216"),
    ("Instrumentación y transacciones", "28", "1064"),
    ("Cuenta corriente", "24", "912"),
    ("Panel de Control", "16", "608"),
    ("Permisos granulares y vistas", "12", "456"),
    ("Experiencia de usuario", "8", "304"),
    ("Responsive y móvil", "12", "456"),
    ("Documentación y despliegue", "12", "456"),
    ("TOTAL", "232", "8816"),
]


def add_presupuesto_slide(prs, presupuesto_data, logo_path=None):
    """Slide con barra de título + tabla de presupuesto y leyenda de fuente."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Inversión – Desarrollo por funcionalidad", logo_path=logo_path)
    # Leyenda clara: tarifa y fuente
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.4))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = "Tarifa de referencia: USD 38/h (Argentina). Fuente: Bitácora del proyecto, solapa Presupuesto."
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRIS_TEXTO
    p2.line_spacing = 1.0
    rows, cols = len(presupuesto_data), 3
    left, top, width = Inches(0.5), Inches(1.32), Inches(12.3)
    row_h = Inches(0.36)
    height = min(row_h * rows, Inches(5.2))
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    for i, fila in enumerate(presupuesto_data):
        for j, val in enumerate(fila[:3]):
            cell = table.cell(i, j)
            cell.text = str(val) if val is not None else ""
            p = cell.text_frame.paragraphs[0]
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = BLANCO
                p.font.size = Pt(11)
                cell.fill.solid()
                cell.fill.fore_color.rgb = AZUL_LOGO
            else:
                p.font.size = Pt(10)
                p.line_spacing = LINE_SPACING_BODY
                p.space_before = Pt(3)
                p.space_after = Pt(3)
            if j == 1 or j == 2:
                p.alignment = PP_ALIGN.RIGHT
    for j in range(3):
        table.cell(rows - 1, j).text_frame.paragraphs[0].font.bold = True
        table.cell(rows - 1, j).fill.solid()
        table.cell(rows - 1, j).fill.fore_color.rgb = RGBColor(236, 253, 245)


def add_pasos_slide(prs, logo_path=None):
    """Próximos pasos con pie de cierre y buen interlineado."""
    pasos = [
        "Reunión de alcance – Validar funcionalidades y prioridades.",
        "Acuerdo comercial – Forma de pago (por hitos o mensual), plazos y aceptación del presupuesto.",
        "Kick-off – Setup del repo, Supabase y primer despliegue.",
        "Desarrollo por módulos – Entregas parciales y pruebas.",
        "Puesta en producción – Dominio, variables de entorno y capacitación.",
    ]
    add_content_slide(prs, "Próximos pasos", pasos, logo_path=logo_path)
    slide = prs.slides[-1]
    box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = "Gracias por considerar esta propuesta. Quedamos a disposición para ajustar alcance o plazos."
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = GRIS_TEXTO
    p.line_spacing = LINE_SPACING_BODY


def add_cierre_slide(prs, logo_path=None):
    """Slide de cierre: barra azul con logo + mensaje de cierre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, prs, "Gracias", logo_path=logo_path)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = "Quedamos a disposición para ajustar alcance o plazos."
    p.font.size = Pt(14)
    p.font.color.rgb = GRIS_TEXTO
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = LINE_SPACING_BODY


def main():
    os.chdir(PROJECT_ROOT)
    logo_path = obtener_logo_png()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    from datetime import datetime
    fecha = datetime.now().strftime("%d/%m/%Y")
    pie = f"Confidencial · {fecha}"

    # 1. Portada
    add_title_slide(prs,
        "Propuesta de desarrollo",
        "Sistema de gestión de órdenes, cajas y cuenta corriente",
        pie,
        logo_path=logo_path
    )

    presupuesto_data = leer_presupuesto_desde_excel() or PRESUPUESTO_FALLBACK

    # 2. El desafío
    add_content_slide(prs, "El desafío – Un solo lugar para operar con claridad", [
        "Operaciones en múltiples monedas (USD, EUR, ARS) y con intermediarios.",
        "Necesidad de trazar cada orden desde la cotización hasta el cobro/pago.",
        "Cajas (efectivo y banco) y cuenta corriente con clientes e intermediarios.",
        "Equipos con distintos niveles de acceso (admin, encargado, visor).",
        "",
        "Centralizar la operación, con seguridad y sin perder el control.",
    ], es_lista=True, logo_path=logo_path)

    # 3. La solución
    add_content_slide(prs, "La solución – Una aplicación web a medida", [
        "Panel de Control con saldos por caja (efectivo/banco), variaciones y pendientes.",
        "Órdenes con estados claros (cotización → cerrada → concertada → ejecutada) e instrumentación.",
        "Cajas con movimientos manuales y automáticos por orden; tipos configurables.",
        "Cuenta corriente con clientes e intermediarios; conciliación y comisiones.",
        "Seguridad por roles y permisos granulares (vistas y acciones).",
        "Responsive para usar en escritorio y móvil.",
        "",
        "Todo en un solo sistema, accesible desde el navegador.",
    ], es_lista=True, logo_path=logo_path)

    # 4. Beneficios
    add_beneficios_slide(prs, logo_path=logo_path)

    # 5. Alcance
    add_alcance_slide(prs, logo_path=logo_path)

    # 6. Presupuesto (desde Bitacora Presupuesto)
    add_presupuesto_slide(prs, presupuesto_data, logo_path=logo_path)

    # 7. Próximos pasos
    add_pasos_slide(prs, logo_path=logo_path)

    # 8. Cierre
    add_cierre_slide(prs, logo_path=logo_path)

    # Paginado: "actual / total" en cada slide (esquina inferior derecha)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        add_slide_number(slide, i + 1, total_slides)

    prs.save(str(OUT_PPTX))
    print(f"Creado: {OUT_PPTX}")
    # Borrar PNG temporal si lo generamos desde SVG
    tmp_logo = SCRIPT_DIR / "_logo.png"
    if tmp_logo.exists():
        try:
            tmp_logo.unlink()
        except Exception:
            pass


if __name__ == "__main__":
    main()
